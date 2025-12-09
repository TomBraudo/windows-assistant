"""
PowerPoint generation tools.

Uses python-pptx to create simple slide decks from structured text, and can
optionally add an image (using the existing image search/downloader).
"""

import os
from typing import List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt

from .os_ops import get_desktop_path
from .image_tools import find_image


def _chunk_list(items: List[str], chunk_size: int) -> List[List[str]]:
    return [items[i : i + chunk_size] for i in range(0, len(items), chunk_size)]


def _extract_saved_path(result: str) -> Optional[str]:
    """
    Given the string result from find_image, extract the actual file path.
    Expected format: 'Saved image to: <path>'.
    """
    prefix = "Saved image to:"
    if not result or prefix not in result:
        return None
    path = result.split(prefix, 1)[1].strip()
    return path or None


def create_presentation(
    title: str,
    bullet_text: str,
    path: str = "presentation.pptx",
    image_query: Optional[str] = None,
) -> str:
    """
    Creates a PowerPoint presentation on the Desktop.

    Args:
        title: Title for the presentation (used on the title slide).
        bullet_text: Newline-separated bullet points. These will be split into
                     multiple slides if there are many.
        path: File path for the .pptx file (relative to Desktop if not absolute).
              Can be just a filename like "My_Presentation.pptx" or a full path.
        image_query: Optional description for an image to add to the final slide.
                     Uses the SerpAPI-based find_image tool to download an image.

    Behavior:
        - Slide 1: Title slide with the given title.
        - Subsequent slides: Content slides with up to 5 bullets per slide.
        - If image_query is provided, the last slide will also include an image
          related to that query (if one can be found and downloaded).
    """
    # Resolve target path (default to Desktop)
    if not os.path.isabs(path):
        path = os.path.join(get_desktop_path(), path)
    os.makedirs(os.path.dirname(path), exist_ok=True)

    prs = Presentation()

    # Title slide
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title
    if slide.placeholders and len(slide.placeholders) > 1:
        try:
            slide.placeholders[1].text = ""
        except Exception:
            pass

    # Prepare bullets
    raw_lines = [line.strip() for line in bullet_text.splitlines()]
    bullets = [line for line in raw_lines if line]

    # Create content slides with up to 5 bullets each
    content_layout = prs.slide_layouts[1]  # Title and Content
    for idx, chunk in enumerate(_chunk_list(bullets, 5), start=1):
        slide = prs.slides.add_slide(content_layout)
        slide_title = slide.shapes.title
        body = slide.placeholders[1]

        slide_title.text = f"{title} - Part {idx}"

        tf = body.text_frame
        tf.clear()
        for i, bullet in enumerate(chunk):
            if i == 0:
                tf.text = bullet
            else:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0

    # Optional image on a final slide
    if image_query:
        image_result = find_image(image_query)
        image_path = _extract_saved_path(image_result)
        if image_path and os.path.exists(image_path):
            img_slide = prs.slides.add_slide(content_layout)
            img_slide.shapes.title.text = f"{title} - Illustration"
            body = img_slide.placeholders[1]
            body.text = image_query

            # Add image roughly centered
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(7)
            try:
                img_slide.shapes.add_picture(image_path, left, top, width=width)
            except Exception:
                # Ignore image insertion errors; the text slide still exists
                pass

    prs.save(path)

    # Try to open the presentation for the user
    try:
        os.startfile(path)
        return f"Created and opened presentation at: {path}"
    except Exception:
        # Even if we can't auto-open, the file is still created
        return f"Created presentation at: {path} (could not auto-open; please open it manually)"


